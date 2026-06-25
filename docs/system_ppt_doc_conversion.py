"""
ForgeSight — submission deck generator.

Produces `ForgeSight_Submission.pptx` (16:9, dark themed). Pipeline:
  1. write 5 PlantUML diagrams to generated_app_diagrams/ and render them to PNG (java -jar plantuml.jar),
  2. build an initials avatar + best-effort fetch the Tata Steel + tech logos (stored under assets/),
  3. assemble the slides, embedding the rendered diagrams and the app screenshots from generated_app_images/.

Run:  python system_ppt_doc_conversion.py
"""
from __future__ import annotations

import json
import subprocess
from pathlib import Path

import requests
from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
DIAG = ROOT / "generated_app_diagrams"
IMG = ROOT / "generated_app_images"
ASSETS = DIAG / "assets"
JAR = DIAG / "plantuml.jar"
OUT = ROOT / "ForgeSight_Submission.pptx"
ASSETS.mkdir(parents=True, exist_ok=True)

# ---- palette (matches the app) -------------------------------------------------
BG = RGBColor(0x0B, 0x0E, 0x13)
PANEL = RGBColor(0x12, 0x17, 0x1F)
CARD = RGBColor(0x16, 0x1B, 0x22)
LINE = RGBColor(0x23, 0x2B, 0x35)
TEXT = RGBColor(0xE6, 0xED, 0xF3)
MUTE = RGBColor(0x8B, 0x98, 0xA5)
ORANGE = RGBColor(0xFF, 0x6A, 0x2B)
BLUE = RGBColor(0x4A, 0x90, 0xD9)
GREEN = RGBColor(0x3F, 0xB6, 0x8B)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
FONT = "Segoe UI"
MONO = "Consolas"

SW, SH = 13.333, 7.5  # inches (16:9)

# ============================================================ PlantUML diagrams
COMMON = (
    "skinparam backgroundColor #FFFFFF\n"
    "skinparam defaultFontName Segoe UI\n"
    "skinparam shadowing false\n"
    "skinparam ArrowColor #4A90D9\n"
    "skinparam rectangleBorderColor #232B35\n"
)

DIAGRAMS: dict[str, str] = {
    "backend_architecture": f"""@startuml
{COMMON}skinparam componentStyle rectangle
skinparam componentBackgroundColor #EEF4FB
skinparam componentBorderColor #4A90D9
skinparam databaseBackgroundColor #FFF1E8
skinparam databaseBorderColor #FF6A2B
title ForgeSight — Backend Architecture
actor Engineer
cloud "Next.js 16 Frontend\\n(Vercel)" as FE
package "FastAPI Backend  (Railway · Docker)" {{
  [REST API\\n/chat /equipment /alerts\\n/work-orders /search /reports] as API
  [Auth\\nSupabase JWT · RLS\\nengineer / admin] as AUTH
  [AgentController\\nLangGraph governed graph] as CTRL
  [Guardrails\\ncitations · LOTO · budgets] as GUARD
  [Tools\\nRAG · ML · text-to-SQL · reports] as TOOLS
  [Scheduler\\nhealth re-scan -> alerts] as SCHED
}}
database "Supabase\\nPostgres + pgvector" as DB
[Synthesis SLM\\nOllama Qwen (on-prem)\\nGroq Llama-3.3 (cloud)] as SLM
[Classical ML\\nanomaly · RUL · failure\\ndefect · Azure-PdM] as ML
Engineer --> FE
FE --> API : HTTPS + Bearer JWT
API --> AUTH
API --> CTRL
CTRL --> GUARD
CTRL --> TOOLS
CTRL --> SLM
TOOLS --> DB
TOOLS --> ML
SCHED --> DB
API --> DB
@enduml
""",

    "rag_pipeline": f"""@startuml
{COMMON}skinparam rectangleBackgroundColor #EEF4FB
title Knowledge-Center Ingestion & RAG Pipeline
rectangle "Equipment Manuals" as M
rectangle "Maintenance SOPs" as S
rectangle "Breakdown / Incident Records" as B
rectangle "Spares & Work Orders" as W
rectangle "Section-aware Chunker" as CH #FFF1E8
rectangle "Embedder\\nnomic-embed-text" as EMB #FFF1E8
database "doc_chunks\\npgvector + tsvector" as VEC
M --> CH
S --> CH
B --> CH
W --> CH
CH --> EMB
EMB --> VEC
actor Engineer
rectangle "Query\\n(fault code / NL question)" as Q
rectangle "Hybrid Retrieval\\nvector + full-text (RRF)" as RET #E9F7F1
rectangle "Enum-constrained\\nCitations" as CIT #E9F7F1
rectangle "SLM Synthesis\\n(grounded, cited card)" as SYN #E9F7F1
Engineer --> Q
Q --> RET
VEC --> RET
RET --> CIT
CIT --> SYN
SYN --> Engineer : cited answer
@enduml
""",

    "multiagent_workflow": f"""@startuml
{COMMON}title Governed Multi-Agent Workflow (LangGraph)
start
:Ingest & Authorize\\n(JWT role · reset per-turn budget);
:Cache lookup;
if (golden cache hit?) then (yes)
  :Serve cached card;
  stop
else (no)
endif
:Classify intent (SLM);
partition "Chartered agents — deterministic tool pipelines" {{
  split
    :Diagnostic\\nRAG + match_history;
  split again
    :Reliability\\nhealth + RUL;
  split again
    :Supervisor\\nscore_priority;
  split again
    :Planner\\nspares + procurement;
  split again
    :Analyst\\ngoverned SQL;
  end split
}}
:Synthesize card\\n(SLM · constrained JSON);
:Guardrail validate\\ncitations · LOTO-first · provenance;
if (COMMIT action?) then (yes)
  :Human gate (approve / reject);
else (no)
endif
:Respond — structured cited card;
stop
@enduml
""",

    "maintenance_lifecycle": f"""@startuml
{COMMON}title Maintenance Decision Lifecycle
state "Fault / Anomaly Alert" as A
state "AI Diagnosis (cited RCA)" as B
state "RUL & Risk Forecast" as C
state "Priority & Bottleneck Scoring" as D
state "Spares & Procurement\\n(lead time vs RUL)" as E
state "Work Order\\n(generate · status · export)" as F
state "Execution\\n(LOTO-first checklist)" as G
state "Feedback & Verification" as H
state "Knowledge Update\\n(verified records)" as I
[*] --> A
A --> B
B --> C
C --> D
D --> E
E --> F
F --> G
G --> H
H --> I
I --> A : continuous improvement
@enduml
""",

    "feedback_loop": f"""@startuml
{COMMON}title Feedback-Driven Improvement Loop (FR-6)
start
:Engineer reviews the cited card;
:Verdict — up / down / "this fixed it";
:POST /feedback;
fork
  :down -> demote the cited record\\n(retrieval re-rank);
fork again
  :fixed -> inject confirmed cause\\n(synthesis few-shot exemplar);
fork again
  :fixed -> breakdown_history.verified = true\\n+ logbook entry;
end fork
:Future answers improve\\n(sharper citations & root-cause);
stop
@enduml
""",
}


def render_diagrams() -> dict[str, Path]:
    out: dict[str, Path] = {}
    for name, src in DIAGRAMS.items():
        puml = DIAG / f"{name}.puml"
        puml.write_text(src, encoding="utf-8")
        try:
            subprocess.run(["java", "-Dfile.encoding=UTF-8", "-jar", str(JAR), "-charset", "UTF-8",
                            "-tpng", str(puml), "-o", str(DIAG)],
                           check=True, capture_output=True, timeout=120)
            png = DIAG / f"{name}.png"
            if png.exists():
                out[name] = png
                print(f"  rendered {png.name}")
            else:
                print(f"  WARN no png for {name}")
        except Exception as e:  # noqa: BLE001
            print(f"  WARN render failed {name}: {str(e)[:120]}")
    return out


# ============================================================ assets (avatar + logos)
def _font(size: int, bold=True) -> ImageFont.FreeTypeFont:
    for cand in (("arialbd.ttf" if bold else "arial.ttf"), "segoeui.ttf", "DejaVuSans-Bold.ttf"):
        try:
            return ImageFont.truetype(cand, size)
        except Exception:  # noqa: BLE001
            continue
    return ImageFont.load_default()


def make_avatar(path: Path, initials="GSG"):
    s = 480
    img = Image.new("RGBA", (s, s), (0, 0, 0, 0))
    d = ImageDraw.Draw(img)
    d.ellipse([8, 8, s - 8, s - 8], fill=(0x16, 0x1B, 0x22, 255), outline=(0xFF, 0x6A, 0x2B, 255), width=10)
    f = _font(190)
    bb = d.textbbox((0, 0), initials, font=f)
    d.text(((s - (bb[2] - bb[0])) / 2 - bb[0], (s - (bb[3] - bb[1])) / 2 - bb[1]), initials,
           font=f, fill=(0xE6, 0xED, 0xF3, 255))
    img.save(path)


def fetch_logo(url: str, path: Path) -> bool:
    try:
        r = requests.get(url, timeout=25, headers={"User-Agent": "Mozilla/5.0"})
        if r.status_code == 200 and len(r.content) > 600 and r.headers.get("content-type", "").startswith("image"):
            path.write_bytes(r.content)
            Image.open(path).verify()
            return True
    except Exception:  # noqa: BLE001
        pass
    return False


def fetch_wikimedia(filename: str, path: Path, width: int = 960) -> bool:
    """Wikimedia renders SVG logos to PNG server-side; grab the thumburl."""
    try:
        api = "https://commons.wikimedia.org/w/api.php"
        p = {"action": "query", "titles": f"File:{filename}", "prop": "imageinfo",
             "iiprop": "url", "iiurlwidth": width, "format": "json"}
        r = requests.get(api, params=p, timeout=25, headers={"User-Agent": "ForgeSight/1.0"}).json()
        for _k, v in r.get("query", {}).get("pages", {}).items():
            ii = v.get("imageinfo")
            if ii:
                url = ii[0].get("thumburl") or ii[0].get("url")
                if url:
                    return fetch_logo(url, path)
    except Exception:  # noqa: BLE001
        pass
    return False


# Tech stack (link stored for every entry; logo fetched best-effort for brand domains).
TECH = {
    "Frontend": [
        ("Next.js 16", "https://nextjs.org", None),
        ("React 19", "https://react.dev", None),
        ("TypeScript", "https://www.typescriptlang.org", None),
        ("Tailwind CSS", "https://tailwindcss.com", "tailwindcss.com"),
        ("react-three-fiber / three.js", "https://threejs.org", "threejs.org"),
        ("Recharts", "https://recharts.org", None),
    ],
    "Backend & Agents": [
        ("FastAPI", "https://fastapi.tiangolo.com", None),
        ("Python 3.12", "https://www.python.org", "python.org"),
        ("LangGraph", "https://www.langchain.com/langgraph", "langchain.com"),
        ("Pydantic", "https://pydantic.dev", "pydantic.dev"),
        ("Uvicorn", "https://www.uvicorn.org", None),
        ("ReportLab (PDF)", "https://www.reportlab.com", None),
    ],
    "AI / ML": [
        ("Qwen2.5-3B via Ollama", "https://ollama.com", "ollama.com"),
        ("Groq · Llama-3.3-70B", "https://groq.com", "groq.com"),
        ("nomic-embed-text", "https://www.nomic.ai", "nomic.ai"),
        ("scikit-learn", "https://scikit-learn.org", None),
        ("XGBoost / LightGBM", "https://xgboost.ai", None),
        ("QLoRA · Unsloth", "https://unsloth.ai", "unsloth.ai"),
    ],
    "Data & Infra": [
        ("Supabase (Postgres+pgvector)", "https://supabase.com", "supabase.com"),
        ("Railway (backend)", "https://railway.app", "railway.app"),
        ("Vercel (frontend)", "https://vercel.com", "vercel.com"),
        ("Docker", "https://www.docker.com", "docker.com"),
        ("GitHub", "https://github.com", "github.com"),
    ],
}


def prepare_assets() -> dict:
    a: dict = {}
    av = ASSETS / "avatar.png"
    make_avatar(av)
    a["avatar"] = av
    tata = ASSETS / "tata_steel.png"
    a["tata"] = tata if fetch_wikimedia("Tata_Steel_Logo.svg", tata, 960) else None
    links = {cat: [{"name": n, "url": u} for (n, u, _d) in items] for cat, items in TECH.items()}
    (ASSETS / "tech_stack_links.json").write_text(json.dumps(links, indent=2), encoding="utf-8")
    a["logos"] = {}
    print(f"  assets ready (tata={'yes' if a['tata'] else 'no'})")
    return a


# ============================================================ slide helpers
prs = Presentation()
prs.slide_width = Inches(SW)
prs.slide_height = Inches(SH)
BLANK = prs.slide_layouts[6]
_n = {"i": 0}


def slide(bg=BG):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s


def box(s, l, t, w, h, fill=None, line=None, line_w=1.0, radius=True):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE,
                             Inches(l), Inches(t), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line; shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def text(s, l, t, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, sp=2):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(sp)
        for (txt, size, color, bold, font) in para:
            r = p.add_run(); r.text = txt
            r.font.size = Pt(size); r.font.color.rgb = color
            r.font.bold = bold; r.font.name = font
    return tb


def accent(s, l, t, w, h, color=ORANGE):
    box(s, l, t, w, h, fill=color, radius=False)


def header(s, kicker, title, sub=None):
    accent(s, 0, 0, 0.16, SH, ORANGE)
    text(s, 0.7, 0.42, 12, 0.4, [[(kicker.upper(), 12, BLUE, True, FONT)]])
    text(s, 0.7, 0.74, 12, 0.9, [[(title, 29, TEXT, True, FONT)]])
    if sub:
        text(s, 0.7, 1.52, 12, 0.5, [[(sub, 13.5, MUTE, False, FONT)]])
    box(s, 0.7, 2.04, 11.93, 0.02, fill=LINE, radius=False)


def footer(s, n):
    text(s, 0.7, 7.08, 8, 0.3, [[("ForgeSight · Tata Steel AI Hackathon 2026 · Gurjas Singh Gandhi", 9, MUTE, False, FONT)]])
    text(s, 11.4, 7.08, 1.2, 0.3, [[(str(n), 9, MUTE, False, FONT)]], align=PP_ALIGN.RIGHT)


def fit(s, path: Path, l, t, w, h, frame=True):
    try:
        iw, ih = Image.open(path).size
    except Exception:  # noqa: BLE001
        box(s, l, t, w, h, fill=CARD, line=LINE)
        text(s, l, t + h / 2 - 0.2, w, 0.4, [[("(image unavailable)", 12, MUTE, False, FONT)]], align=PP_ALIGN.CENTER)
        return
    ar, bar = iw / ih, w / h
    if ar > bar:
        nw = w; nh = w / ar
    else:
        nh = h; nw = h * ar
    nl = l + (w - nw) / 2; nt = t + (h - nh) / 2
    if frame:
        box(s, nl - 0.05, nt - 0.05, nw + 0.1, nh + 0.1, fill=CARD, line=LINE)
    s.shapes.add_picture(str(path), Inches(nl), Inches(nt), Inches(nw), Inches(nh))


def bullets(s, items, l=0.95, t=2.35, w=11.5, h=4.4, size=16, gap=8):
    paras = []
    for it in items:
        if isinstance(it, tuple):
            paras.append([("▸  ", 16, ORANGE, True, FONT), (it[0], size, TEXT, True, FONT),
                          (it[1], size, MUTE, False, FONT)])
        else:
            paras.append([("▸  ", 16, ORANGE, True, FONT), (it, size, TEXT, False, FONT)])
    text(s, l, t, w, h, paras, sp=gap)


def page():
    _n["i"] += 1
    return _n["i"]


# ============================================================ build deck
def build(diagrams: dict, assets: dict):
    # 1 · TITLE
    s = slide()
    accent(s, 0, 0, SW, 0.22, ORANGE)
    accent(s, 0, SH - 0.22, SW, 0.22, BLUE)
    if assets.get("tata"):
        box(s, 0.7, 0.55, 3.0, 1.1, fill=WHITE)
        fit(s, assets["tata"], 0.85, 0.66, 2.7, 0.88, frame=False)
    else:
        text(s, 0.7, 0.7, 6, 0.8, [[("TATA STEEL", 26, RGBColor(0x6A, 0xA9, 0xE0), True, FONT)]])
    text(s, 0.7, 2.35, 12, 1.1, [[("Forge", 60, TEXT, True, FONT), ("Sight", 60, ORANGE, True, FONT)]])
    text(s, 0.72, 3.55, 12, 0.6, [[("Intelligent Maintenance Wizard for Steel Plants", 22, TEXT, False, FONT)]])
    text(s, 0.72, 4.15, 12, 0.5, [[("Governed multi-agent, citation-grounded maintenance decision-support", 13.5, MUTE, False, FONT)]])
    text(s, 0.72, 4.6, 12, 0.4, [[("Tata Steel AI Hackathon 2026", 14, BLUE, True, FONT)]])
    box(s, 8.6, 4.9, 4.05, 1.95, fill=CARD, line=LINE)
    fit(s, assets["avatar"], 8.85, 5.15, 1.45, 1.45, frame=False)
    text(s, 10.45, 5.2, 2.1, 1.7, [
        [("Gurjas Singh Gandhi", 15, TEXT, True, FONT)],
        [("Software Engineer", 12, ORANGE, False, FONT)],
        [("Pune, Maharashtra, India", 11, MUTE, False, FONT)],
    ], sp=4)
    footer(s, page())

    # 2 · PROBLEM
    s = slide(); header(s, "Problem Statement", "Maintenance in steel plants is fragmented & reactive")
    bullets(s, [
        ("Capital-intensive, interdependent equipment — ", "any unplanned downtime cascades into production loss, safety risk and cost."),
        ("Knowledge is scattered — ", "manuals, SOPs, breakdown logs, failure reports and sensor alerts live in silos."),
        ("Diagnosis is manual & expert-dependent — ", "slow response, inconsistent decisions, hard to onboard."),
        ("No single context-aware support — ", "engineers juggle sources to find root cause and the right corrective action."),
        ("The need — ", "an explainable system that consolidates evidence, diagnoses, predicts degradation, prioritises actions and learns."),
    ])
    footer(s, page())

    # 3 · IDEA & OBJECTIVE
    s = slide(); header(s, "Idea & Objective", "A governed AI copilot: fault code → auditable fix plan in ~90s")
    bullets(s, [
        ("Consolidate evidence — ", "RAG over manuals/SOPs/records + governed SQL over operational data."),
        ("Diagnose & explain — ", "ranked root causes, every claim cited to a manual, record or trend (cite-or-refuse)."),
        ("Predict — ", "anomaly detection, RUL forecast and early warning for catastrophic failure."),
        ("Prioritise — ", "deterministic risk score from criticality, delay severity, spares & lead time."),
        ("Act & learn — ", "work orders, HITL approval, and a feedback loop that improves future answers."),
        ("On-prem capable — ", "fine-tuned open-source SLM; no plant data needs to leave the network."),
    ], size=15, gap=7)
    footer(s, page())

    # 4 · INNOVATION & IMPACT
    s = slide(); header(s, "Innovation & Business Impact", "Trust is the differentiator")
    col = [
        ("Governed multi-agent — ", "5 chartered agents run deterministic tool pipelines (not free ReAct); every action audited."),
        ("Citations-or-refuse — ", "a code-level guardrail makes a fabricated citation structurally impossible."),
        ("Deterministic scoring — ", "priority & procurement from auditable rules, never an LLM guess."),
        ("Human-in-the-loop — ", "agents propose; engineers commit; every allow/deny timestamped."),
    ]
    impact = [
        ("Reduce unplanned downtime", "earlier, RUL-driven intervention"),
        ("Faster response", "fault → cited fix plan in ~90 seconds"),
        ("Higher diagnostic accuracy", "evidence-grounded, consistent"),
        ("Proactive maintenance", "reactive → predictive shift"),
        ("Better spares planning", "lead time vs RUL exposure"),
        ("Leadership ROI view", "shutdown vs failure cost, savings"),
    ]
    bullets(s, col, l=0.95, t=2.35, w=6.0, h=4.4, size=14, gap=8)
    box(s, 7.3, 2.35, 5.3, 4.35, fill=PANEL, line=LINE)
    text(s, 7.55, 2.5, 5, 0.4, [[("BUSINESS IMPACT", 12, BLUE, True, FONT)]])
    paras = [[("•  ", 13, GREEN, True, FONT), (a + " — ", 13, TEXT, True, FONT), (b, 13, MUTE, False, FONT)] for a, b in impact]
    text(s, 7.55, 2.95, 5.0, 3.6, paras, sp=9)
    footer(s, page())

    # 5 · ARCHITECTURE
    s = slide(); header(s, "System Architecture", "Three governed layers over a vector-enabled store")
    bullets(s, [
        ("Frontend (Next.js 16 · Vercel) — ", "landing, role-based auth, dashboard tabs, equipment console, 3D digital twin."),
        ("Backend (FastAPI · Railway · Docker) — ", "LangGraph AgentController, Authority, Guardrails, tools, scheduler."),
        ("Data (Supabase · Postgres + pgvector) — ", "equipment, sensors, health, alerts, breakdowns, doc_chunks, spares, work_orders, RLS."),
        ("Reasoning — ", "SLM synthesis (Ollama Qwen on-prem / Groq Llama-3.3 cloud) under constrained JSON decoding."),
        ("Trust tiers — ", "Tier-1 schema+RLS · Tier-2 charter governance + audit log · Tier-3 output guardrails."),
    ], size=15, gap=9)
    footer(s, page())

    # 6 · DATA & SYSTEM FLOW
    s = slide(); header(s, "Data Flow & System Flow", "From sensor to cited decision")
    bullets(s, [
        ("Ingest — ", "physics-shaped sensor stream + corpus (manuals/SOPs/records) embedded into doc_chunks."),
        ("Detect — ", "scheduler re-scans health (anomaly + RUL) → raises severity-ranked alerts."),
        ("Ask — ", "engineer query / alert enters the governed graph: authorize → cache → classify intent."),
        ("Reason — ", "chartered agent runs its deterministic tool pipeline (RAG · ML · SQL) gathering evidence."),
        ("Synthesize — ", "SLM fills a typed card; guardrails verify citations exist & LOTO-first ordering."),
        ("Act — ", "COMMIT actions pause for human approval; outcome captured as feedback → improves retrieval."),
    ], size=14, gap=7)
    footer(s, page())

    # 7 · MODEL DESIGN
    s = slide(); header(s, "Model Design & Reasoning Pipeline", "SLM narrates; tools decide; numbers are deterministic")
    bullets(s, [
        ("Classical ML (live) — ", "IsolationForest+EWMA anomaly, XGBoost RUL (C-MAPSS), failure (AI4I), LightGBM defect, Azure-PdM 24h."),
        ("Retrieval — ", "hybrid vector + full-text (RRF), metadata-filtered; cloud degrades to full-text-primary, citations stay real."),
        ("SLM role — ", "invoked only at synthesize/repair; constrained JSON decoding; never selects tools or computes numbers."),
        ("Fine-tune — ", "QLoRA on Qwen2.5-3B (Unsloth, Colab T4) → GGUF → Ollama; promotion gated vs base; base is the safe fallback."),
        ("Guardrails — ", "citation-existence, LOTO-first, matrix-provenance, SELECT-only SQL — enforced in code, not prompts."),
    ], size=14, gap=8)
    footer(s, page())

    # 8-12 · DIAGRAMS
    diag_meta = [
        ("backend_architecture", "Architecture", "Backend Architecture",
         "FastAPI governed graph over Supabase (pgvector); SLM synthesis with on-prem/cloud serving; classical ML tools."),
        ("rag_pipeline", "Knowledge Center", "RAG / Knowledge Ingestion Pipeline",
         "Manuals, SOPs, incidents, spares & work orders are chunked, embedded and hybrid-retrieved as cited evidence."),
        ("multiagent_workflow", "Orchestration", "Multi-Agent Workflow",
         "Authorize → cache → classify → chartered agent pipelines → synthesize → guardrails → human gate → cited card."),
        ("maintenance_lifecycle", "Lifecycle", "Maintenance Decision Lifecycle",
         "Alert → diagnosis → RUL → priority → spares → work order → execution → feedback → knowledge update."),
        ("feedback_loop", "Learning", "Feedback-Driven Improvement Loop",
         "up/down/fixed verdicts re-rank retrieval, inject verified exemplars, and mark records engineer-verified."),
    ]
    for key, kicker, title, desc in diag_meta:
        s = slide(); header(s, kicker, title, desc)
        if key in diagrams:
            fit(s, diagrams[key], 0.9, 2.25, 11.5, 4.6)
        else:
            box(s, 0.9, 2.25, 11.5, 4.6, fill=CARD, line=LINE)
            text(s, 0.9, 4.3, 11.5, 0.5, [[("(diagram render unavailable)", 14, MUTE, False, FONT)]], align=PP_ALIGN.CENTER)
        footer(s, page())

    # 13 · TECH STACK
    s = slide(); header(s, "Technology Stack", "Open-source first, deploy-ready")
    cw, ch, cols = 5.85, 2.05, 2
    for i, (cat, items) in enumerate(TECH.items()):
        cx = 0.8 + (i % cols) * (cw + 0.25)
        cy = 2.35 + (i // cols) * (ch + 0.22)
        box(s, cx, cy, cw, ch, fill=PANEL, line=LINE)
        text(s, cx + 0.2, cy + 0.12, cw - 0.4, 0.35, [[(cat.upper(), 12, BLUE, True, FONT)]])
        paras = [[("•  ", 12, ORANGE, True, FONT), (name, 12, TEXT, False, FONT)] for (name, _u, _d) in items]
        text(s, cx + 0.2, cy + 0.52, cw - 0.4, ch - 0.6, paras, sp=3)
    text(s, 0.8, 6.66, 11.5, 0.3,
         [[("Official links stored in generated_app_diagrams/assets/tech_stack_links.json", 10, MUTE, False, FONT)]])
    footer(s, page())

    # 14+ · SCREENSHOTS
    shots = [
        ("home.png", "Product", "Public landing — fault code to cited fix plan", "App-flow explainer, governed multi-agent + on-prem SLM pitch, the three scenarios."),
        ("login.png", "Access", "Authentication — Supabase JWT, role-based", "Engineer & admin roles, verified on every request; engineer-only public signup."),
        ("dashboard.png", "Console", "Plant Overview — computed KPIs & live health", "Criticality-weighted availability, assets-alerting, downtime-at-risk; tiles + alert feed."),
        ("twin-3d.png", "Digital Twin", "3D plant twin — zones & health", "Assets shaped by type, color-coded by health; orbit, zoom, click to inspect."),
        ("twin-3d-inspect.png", "Digital Twin", "Twin inspector — RUL, status, work orders", "Per-asset anomaly, RUL, maintenance status and open work orders from live data."),
        ("equipment-sinter-fan-2.png", "Operational Console", "Asset console — live risk + copilot", "Health, RUL band, sensor trend and the governed copilot in one screen."),
        ("equipment-sinter-fan-2-diagnosis.png", "AI Investigation", "AI-guided diagnosis with Evidence trail", "Ranked root causes, confidence, citations resolvable in the Evidence drawer."),
        ("agent-diagnostic-diagnosis.png", "Agent", "Diagnostic agent — cited root-cause analysis", "RAG + verified breakdown records; cite-or-refuse."),
        ("agent-diagnostic-checklist.png", "Agent", "Checklist / SOP — LOTO-first repair steps", "Safety-highlighted procedure, citations to the exact SOP."),
        ("agent-reliability-rul.png", "Agent", "Reliability agent — RUL estimate", "Trend extrapolation validated by the C-MAPSS XGBoost method."),
        ("agent-reliability-wait.png", "Agent", "Wait-assessment — multi-agent fan-out", "Can it wait till Sunday? Reliability + Planner + Supervisor + monitoring plan."),
        ("agent-supervisor-priority.png", "Agent", "Supervisor — deterministic priority score", "Factor breakdown (criticality/delay/spares/lead); never an LLM guess."),
        ("agent-planner-spares.png", "Agent", "Planner — spares & procurement", "Stock, lead time vs RUL, reserve/PO proposal under HITL."),
        ("agent-analyst-sql.png", "Agent", "Analyst — governed text-to-SQL", "SELECT-only over curated views; the SQL is the citation."),
        ("dashboard-evidence.png", "Module", "Evidence search — unified knowledge", "Manuals, SOPs, incidents, sensor events, work orders & spares as searchable evidence."),
        ("dashboard-work-orders.png", "Module", "Work orders — execution flow", "Status tracking, JSON/PDF export, maintenance execution."),
        ("dashboard-incidents.png", "Module", "Incident replay & lessons learned", "Production impact, failure progression, corrective action, similar-failure lessons."),
        ("dashboard-spares.png", "Module", "Spare catalog + inventory optimizer", "Stock, cost, lead time, linked asset, PO action; shortage risk & production exposure."),
        ("dashboard-reliability.png", "Module", "Reliability analytics", "Predictive curves, failure probability, RUL forecast, trend analysis."),
        ("dashboard-leadership.png", "Module", "Leadership ROI review", "Shutdown vs potential-failure cost, expected savings, ROI, confidence, recommended action."),
    ]
    for fn, kicker, title, desc in shots:
        p = IMG / fn
        if not p.exists():
            continue
        s = slide(); header(s, kicker, title, desc)
        fit(s, p, 0.9, 2.25, 11.5, 4.6)
        footer(s, page())

    # ALERTING & PREDICTION
    s = slide(); header(s, "Alerting & Prediction Logic", "Real-time abnormality detection & early warning")
    bullets(s, [
        ("Anomaly — ", "IsolationForest on standardized features + EWMA control limits; recall-first for early warning."),
        ("Severity rule — ", "sustained-window confirmation; CRITICAL when RUL < spares lead time (catastrophic-risk gate)."),
        ("RUL — ", "robust trend extrapolation to the next vibration limit; XGBoost C-MAPSS validates the method."),
        ("Failure / defect / PdM — ", "XGBoost & LightGBM benchmark models surfaced as live held-out inferences."),
        ("Real-time — ", "FastAPI lifespan scheduler re-scans every interval → alerts table → /alerts feed + UI toasts."),
        ("Explainable — ", "every alert carries contributing sensors and resolvable evidence."),
    ], size=14, gap=7)
    footer(s, page())

    # ASSUMPTIONS & LIMITATIONS
    s = slide(); header(s, "Assumptions & Limitations", "Honest framing")
    bullets(s, [
        ("Digital twin — ", "sensor stream is a physics-shaped simulation (no public steel-plant feed); governance/ML/reasoning are real."),
        ("Benchmark models — ", "failure/RUL/defect/PdM validate the method on public datasets; not per-asset sensor models."),
        ("Cloud serving — ", "public demo runs Groq (no GPU) + full-text-primary retrieval; on-prem runs fine-tuned Qwen via Ollama."),
        ("Fine-tune — ", "QLoRA pipeline reproducible on Colab T4; base Qwen ships as the gated, safe fallback."),
        ("Costs / ROI — ", "downtime-at-risk uses a documented ₹/hr × criticality assumption, returned with every figure."),
    ], size=14, gap=8)
    footer(s, page())

    # INSTALL / RUN
    s = slide(); header(s, "Install · Configure · Run", "Deployed & reproducible")
    box(s, 0.8, 2.35, 11.8, 1.55, fill=PANEL, line=LINE)
    text(s, 1.05, 2.5, 11.3, 1.35, [
        [("Live frontend (Vercel):  ", 14, TEXT, True, FONT), ("https://forge-sight-one.vercel.app", 14, BLUE, False, MONO)],
        [("Backend (Railway):  ", 14, TEXT, True, FONT), ("railway.com/project/7db214e3-1857-4961-947d-cd1ecc851d3f", 12.5, BLUE, False, MONO)],
        [("Source (GitHub):  ", 14, TEXT, True, FONT), ("https://github.com/Gurjas2112/ForgeSight", 14, BLUE, False, MONO)],
        [("Demo login:  ", 14, TEXT, True, FONT), ("engineer@demo.forgesight / forgesight-demo", 12.5, MUTE, False, MONO)],
    ], sp=8)
    text(s, 0.95, 4.18, 11.5, 0.4, [[("LOCAL RUN", 12, BLUE, True, FONT)]])
    code = [
        "cp .env.example .env            # DATABASE_URL + keys",
        "uv sync                          # backend + ml deps",
        "ollama pull qwen2.5:3b-instruct && ollama pull nomic-embed-text",
        "python backend/db/apply_migrations.py     # schema + seeds + corpus",
        "uvicorn backend.server:app --port 8000    # backend",
        "cd frontend && npm install && npm run dev # http://localhost:3000",
    ]
    box(s, 0.8, 4.52, 11.8, 2.2, fill=RGBColor(0x0E, 0x11, 0x16), line=LINE)
    text(s, 1.05, 4.67, 11.3, 2.0, [[(c, 12, GREEN, False, MONO)] for c in code], sp=6)
    footer(s, page())

    # THANK YOU
    s = slide()
    accent(s, 0, 0, SW, 0.22, ORANGE)
    accent(s, 0, SH - 0.22, SW, 0.22, BLUE)
    text(s, 0, 2.7, SW, 1.0, [[("Thank you", 54, TEXT, True, FONT)]], align=PP_ALIGN.CENTER)
    text(s, 0, 3.85, SW, 0.5, [[("ForgeSight — governed, explainable maintenance intelligence for steel plants", 16, MUTE, False, FONT)]], align=PP_ALIGN.CENTER)
    text(s, 0, 4.55, SW, 0.5, [[("Gurjas Singh Gandhi · Software Engineer · Pune, India", 14, BLUE, True, FONT)]], align=PP_ALIGN.CENTER)
    text(s, 0, 5.05, SW, 0.4, [[("forge-sight-one.vercel.app  ·  github.com/Gurjas2112/ForgeSight", 12, MUTE, False, MONO)]], align=PP_ALIGN.CENTER)
    page()

    prs.save(OUT)
    print(f"\nSaved {OUT.name} · {len(prs.slides._sldIdLst)} slides")


def main():
    print("1/3 rendering diagrams…")
    diagrams = render_diagrams()
    print("2/3 preparing assets…")
    assets = prepare_assets()
    print("3/3 building deck…")
    build(diagrams, assets)


if __name__ == "__main__":
    main()
